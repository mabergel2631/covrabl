"""Generate the Covrabl pitch deck — 8 slides + optional closing.

Mirrors `DECK_DRAFT.md` exactly. On-stage copy only; no speaker notes
embedded (per the user's instruction — slide deck stays clean).

Output: %DESKTOP%\\Covrabl-Deck.pptx
"""
from __future__ import annotations

import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu


# ── Output location ─────────────────────────────────────────────────────────
# Windows Desktop is sometimes redirected to OneDrive, sometimes not. To be
# robust, save to BOTH locations if they exist — whichever one the user's
# File Explorer "Desktop" view points at, they'll see the file.
_onedrive_desktop = Path(r"C:\Users\maber\OneDrive\Desktop")
_bare_desktop = Path(r"C:\Users\maber\Desktop")
OUTPUT_PATHS = [d / "Covrabl-Deck.pptx" for d in (_onedrive_desktop, _bare_desktop) if d.exists()]
OUTPUT_PATH = OUTPUT_PATHS[0] if OUTPUT_PATHS else _bare_desktop / "Covrabl-Deck.pptx"

CONTACT_EMAIL = "support@covrabl.com"
WEBSITE = "covrabl.com"


# ── Brand palette ───────────────────────────────────────────────────────────
NAVY_DARK = RGBColor(0x0F, 0x1F, 0x33)      # hero background
NAVY = RGBColor(0x18, 0x2E, 0x4C)
PRIMARY = RGBColor(0x14, 0x4F, 0x6B)         # var(--color-primary) approx
PRIMARY_DARK = RGBColor(0x0B, 0x39, 0x52)
SECONDARY = RGBColor(0x3F, 0xA7, 0xA3)       # teal accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF8, 0xFA, 0xFC)
TEXT = RGBColor(0x0F, 0x17, 0x2A)
TEXT_MUTED = RGBColor(0x64, 0x74, 0x8B)
TEXT_SECONDARY = RGBColor(0x47, 0x55, 0x69)
SURFACE = RGBColor(0xF1, 0xF5, 0xF9)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
ACCENT_GREEN = RGBColor(0x16, 0x65, 0x34)
ACCENT_AMBER = RGBColor(0x92, 0x40, 0x0E)
ACCENT_RED = RGBColor(0xDC, 0x26, 0x26)
ACCENT_CYAN = RGBColor(0x08, 0x91, 0xB2)
ACCENT_SLATE = RGBColor(0x64, 0x74, 0x8B)


# ── Slide geometry — 16:9 widescreen at 13.333" × 7.5" ──────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── Helpers ─────────────────────────────────────────────────────────────────

def add_blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(layout)
    return slide


def set_slide_bg(slide, color: RGBColor):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill: RGBColor | None = None,
             line_color: RGBColor | None = None, line_width: float = 1.0):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    return shape


def add_text(slide, left, top, width, height, text: str, *,
             font_size: int = 18, bold: bool = False, color: RGBColor = TEXT,
             align: int = PP_ALIGN.LEFT, anchor: int = MSO_ANCHOR.TOP,
             italic: bool = False, font_name: str = "Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_multi(slide, left, top, width, height, paragraphs: list[dict], *,
              anchor: int = MSO_ANCHOR.TOP):
    """Add a textbox with multiple paragraphs.
    Each paragraph is a dict: { text, font_size, bold, color, align, space_after, italic }.
    """
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = para.get("align", PP_ALIGN.LEFT)
        if "space_after" in para:
            p.space_after = Pt(para["space_after"])
        run = p.add_run()
        run.text = para["text"]
        run.font.name = "Calibri"
        run.font.size = Pt(para.get("font_size", 18))
        run.font.bold = para.get("bold", False)
        run.font.italic = para.get("italic", False)
        run.font.color.rgb = para.get("color", TEXT)
    return tb


def page_footer(slide, page_num: int, dark: bool = False):
    """Small page number + brand tag at the bottom of every slide."""
    color = RGBColor(0xCB, 0xD5, 0xE1) if dark else TEXT_MUTED
    add_text(slide, Inches(0.5), Inches(7.05), Inches(2.0), Inches(0.3),
             f"COVRABL  ·  {page_num}", font_size=9, color=color)
    add_text(slide, Inches(10.833), Inches(7.05), Inches(2.0), Inches(0.3),
             WEBSITE, font_size=9, color=color, align=PP_ALIGN.RIGHT)


# ── Slide 1: Between Renewals (hero) ────────────────────────────────────────

def slide_1(prs: Presentation):
    s = add_blank_slide(prs)
    set_slide_bg(s, NAVY_DARK)

    # Headline — large white
    add_text(s, Inches(1.0), Inches(2.6), Inches(11.333), Inches(1.6),
             "Insurance relationships shouldn't disappear between renewals.",
             font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

    # Subline
    add_text(s, Inches(1.0), Inches(4.4), Inches(11.333), Inches(0.5),
             "Today, most do.",
             font_size=22, color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.CENTER, italic=True)

    # Timeline visual — calendar dots showing 11 silent months between two renewals
    timeline_y = Inches(5.6)
    add_rect(s, Inches(2.0), timeline_y, Inches(9.333), Emu(1500),
             fill=RGBColor(0x47, 0x55, 0x69))
    # Renewal pins at start and end
    for x in [Inches(2.0), Inches(11.333)]:
        pin = s.shapes.add_shape(MSO_SHAPE.OVAL, x - Inches(0.12),
                                  timeline_y - Inches(0.1), Inches(0.24), Inches(0.24))
        pin.fill.solid()
        pin.fill.fore_color.rgb = SECONDARY
        pin.line.fill.background()
    # Labels
    add_text(s, Inches(1.5), Inches(5.95), Inches(1.5), Inches(0.3),
             "RENEWAL", font_size=9, bold=True, color=RGBColor(0xCB, 0xD5, 0xE1),
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(10.8), Inches(5.95), Inches(1.5), Inches(0.3),
             "RENEWAL", font_size=9, bold=True, color=RGBColor(0xCB, 0xD5, 0xE1),
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(2.0), Inches(5.95), Inches(9.333), Inches(0.3),
             "11 months of silence",
             font_size=11, italic=True, color=RGBColor(0x94, 0xA3, 0xB8),
             align=PP_ALIGN.CENTER)

    page_footer(s, 1, dark=True)


# ── Slide 2: Covrabl Overview ───────────────────────────────────────────────

def slide_2(prs: Presentation):
    s = add_blank_slide(prs)
    set_slide_bg(s, WHITE)

    # Eyebrow
    add_text(s, Inches(0.8), Inches(0.6), Inches(12.0), Inches(0.4),
             "WHAT COVRABL DOES", font_size=10, bold=True, color=PRIMARY)

    # Headline
    add_text(s, Inches(0.8), Inches(1.05), Inches(12.0), Inches(1.0),
             "A shared workspace where clients and their agents stay connected to coverage — between every renewal.",
             font_size=24, bold=True, color=TEXT)

    # Tagline
    add_text(s, Inches(0.8), Inches(2.3), Inches(12.0), Inches(0.9),
             ("Covrabl uses AI to help clients understand their coverage and agents spot "
              "what's worth their time — between every renewal."),
             font_size=14, color=TEXT_SECONDARY)

    # Two pillar cards
    card_top = Inches(3.6)
    card_h = Inches(2.7)
    card_w = Inches(5.8)
    gap = Inches(0.4)
    left1 = Inches(0.8)
    left2 = left1 + card_w + gap

    for left, eyebrow, title, line in [
        (left1, "PILLAR 1",
         "Identify the conversation worth having",
         "The right client. The right reason. The right week."),
        (left2, "PILLAR 2",
         "Have the conversation well",
         "Show your work. Make it shareable. Make it land."),
    ]:
        add_rect(s, left, card_top, card_w, card_h, fill=SURFACE,
                 line_color=BORDER, line_width=1)
        add_text(s, left + Inches(0.4), card_top + Inches(0.35),
                 card_w - Inches(0.8), Inches(0.3),
                 eyebrow, font_size=10, bold=True, color=PRIMARY)
        add_text(s, left + Inches(0.4), card_top + Inches(0.75),
                 card_w - Inches(0.8), Inches(1.0),
                 title, font_size=20, bold=True, color=TEXT)
        add_text(s, left + Inches(0.4), card_top + Inches(1.85),
                 card_w - Inches(0.8), Inches(0.8),
                 line, font_size=14, color=TEXT_SECONDARY, italic=True)

    # Bottom tag
    add_text(s, Inches(0.8), Inches(6.6), Inches(12.0), Inches(0.3),
             "Shared insurance intelligence between clients and their trusted advisors.",
             font_size=11, color=TEXT_MUTED, italic=True, align=PP_ALIGN.CENTER)

    page_footer(s, 2)


# ── Slide 3: What Clients See ───────────────────────────────────────────────

def slide_3(prs: Presentation):
    s = add_blank_slide(prs)
    set_slide_bg(s, WHITE)

    add_text(s, Inches(0.8), Inches(0.6), Inches(12.0), Inches(0.4),
             "FOR YOUR BOOK", font_size=10, bold=True, color=PRIMARY)
    add_text(s, Inches(0.8), Inches(1.05), Inches(12.0), Inches(1.0),
             "This is what your book sees.",
             font_size=30, bold=True, color=TEXT)

    # Two-column layout — copy on left, mock portal on right
    add_multi(s, Inches(0.8), Inches(2.4), Inches(5.8), Inches(4.0), [
        {"text": "A clean, branded view of every policy you've shared with them.",
         "font_size": 16, "color": TEXT, "bold": True, "space_after": 12},
        {"text": ("No ads. No quote spam. No \"we noticed you might also be "
                  "interested in life insurance.\" Just their coverage — "
                  "organized, current, theirs to reference any time."),
         "font_size": 13, "color": TEXT_SECONDARY, "space_after": 14},
        {"text": "✓  Branded as your agency, not Covrabl",
         "font_size": 13, "color": TEXT_SECONDARY, "space_after": 6},
        {"text": "✓  No data sold. No marketing emails. Not now, not ever.",
         "font_size": 13, "color": TEXT_SECONDARY, "space_after": 6},
        {"text": "✓  Renewal reviews and quote comparisons land here too",
         "font_size": 13, "color": TEXT_SECONDARY},
    ])

    # Right side — mock client portal
    portal_left = Inches(7.2)
    portal_top = Inches(2.4)
    portal_w = Inches(5.3)
    portal_h = Inches(4.0)
    add_rect(s, portal_left, portal_top, portal_w, portal_h,
             fill=WHITE, line_color=BORDER, line_width=1)
    # Portal header bar
    add_rect(s, portal_left, portal_top, portal_w, Inches(0.6),
             fill=SURFACE, line_color=BORDER, line_width=1)
    add_text(s, portal_left + Inches(0.2), portal_top + Inches(0.1),
             Inches(3.0), Inches(0.3),
             "Your Coverage", font_size=11, bold=True, color=TEXT)
    add_text(s, portal_left + Inches(0.2), portal_top + Inches(0.32),
             Inches(3.0), Inches(0.25),
             "3 policies on file", font_size=9, color=TEXT_MUTED)
    add_text(s, portal_left + Inches(2.5), portal_top + Inches(0.18),
             portal_w - Inches(2.7), Inches(0.3),
             "Branded: Westlake Insurance Agency",
             font_size=9, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)

    # Three policy rows
    rows = [
        ("Allstate Auto", "Renews Oct 3, 2026", "$1,980/yr"),
        ("Allstate Home", "Renews Oct 3, 2026", "$2,900/yr"),
        ("Allstate Umbrella", "Renews Oct 3, 2026", "$480/yr"),
    ]
    row_top = portal_top + Inches(0.7)
    row_h = Inches(0.85)
    for i, (carrier, renew, premium) in enumerate(rows):
        y = row_top + Inches(i * 0.95)
        add_rect(s, portal_left, y, portal_w, row_h,
                 fill=WHITE, line_color=BORDER, line_width=0.5)
        add_text(s, portal_left + Inches(0.2), y + Inches(0.15),
                 Inches(3.0), Inches(0.3),
                 carrier, font_size=12, bold=True, color=TEXT)
        add_text(s, portal_left + Inches(0.2), y + Inches(0.45),
                 Inches(3.0), Inches(0.3),
                 renew, font_size=9, color=TEXT_MUTED)
        add_text(s, portal_left + Inches(3.5), y + Inches(0.25),
                 Inches(1.6), Inches(0.4),
                 premium, font_size=11, bold=True, color=TEXT,
                 align=PP_ALIGN.RIGHT)

    add_text(s, Inches(7.2), Inches(6.55), Inches(5.3), Inches(0.3),
             "Branded as your agency. Not Covrabl.",
             font_size=10, color=TEXT_MUTED, italic=True, align=PP_ALIGN.CENTER)

    page_footer(s, 3)


# ── Slide 4: What Clients Can Do — the workspace ────────────────────────────

def slide_4(prs: Presentation):
    s = add_blank_slide(prs)
    set_slide_bg(s, WHITE)

    add_text(s, Inches(0.8), Inches(0.6), Inches(12.0), Inches(0.4),
             "WHAT YOUR CLIENTS CAN DO", font_size=10, bold=True, color=PRIMARY)
    add_text(s, Inches(0.8), Inches(1.05), Inches(12.0), Inches(1.0),
             "It's a workspace, not just a portal.",
             font_size=30, bold=True, color=TEXT)
    add_text(s, Inches(0.8), Inches(2.05), Inches(12.0), Inches(0.6),
             ("The same branded view that surfaces your renewal reviews is "
              "also where your book runs the rest of their insurance life."),
             font_size=14, color=TEXT_SECONDARY)

    # 2 x 2 grid of capability tiles
    tiles = [
        ("💬", "Ask Covrabl, in plain English",
         "Conversational AI explains any policy term, exclusion, or limit. No jargon, no scavenger hunts."),
        ("📜", "Verify against requirements",
         "Upload a lease, loan, or vendor contract. We compare the actual policy against what's required — not just the COI."),
        ("🚨", "Emergency access",
         "One-tap SOS card with carrier hotlines, claim numbers, and key coverage limits. Available without login."),
        ("👨‍👩‍👧", "Share with the household",
         "Family, partner, executor sharing. The people who'd need access in a crisis already have it."),
    ]
    tile_top = Inches(2.95)
    tile_h = Inches(1.7)
    tile_w = Inches(5.95)
    gap_x = Inches(0.25)
    gap_y = Inches(0.2)
    left0 = Inches(0.8)

    for i, (icon, title, body) in enumerate(tiles):
        col = i % 2
        row = i // 2
        left = left0 + col * (tile_w + gap_x)
        top = tile_top + row * (tile_h + gap_y)
        add_rect(s, left, top, tile_w, tile_h, fill=SURFACE,
                 line_color=BORDER, line_width=1)
        # Icon
        add_text(s, left + Inches(0.35), top + Inches(0.3),
                 Inches(0.8), Inches(0.8),
                 icon, font_size=28, color=PRIMARY)
        # Title
        add_text(s, left + Inches(1.2), top + Inches(0.32),
                 tile_w - Inches(1.5), Inches(0.5),
                 title, font_size=15, bold=True, color=TEXT)
        # Body
        add_text(s, left + Inches(1.2), top + Inches(0.85),
                 tile_w - Inches(1.5), Inches(0.75),
                 body, font_size=11, color=TEXT_SECONDARY)

    # Footer line
    add_text(s, Inches(0.8), Inches(6.85), Inches(12.0), Inches(0.3),
             ("Every one of these is something your agency made available to them. "
              "They don't see Covrabl marketing — they see what their agent gave them."),
             font_size=10, italic=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    page_footer(s, 4)


# ── Slide 5: What Agents See — This Week (hero) ─────────────────────────────

def slide_5(prs: Presentation):
    s = add_blank_slide(prs)
    set_slide_bg(s, OFF_WHITE)

    add_text(s, Inches(0.8), Inches(0.6), Inches(12.0), Inches(0.4),
             "FOR THE AGENT", font_size=10, bold=True, color=PRIMARY)
    add_text(s, Inches(0.8), Inches(1.05), Inches(12.0), Inches(0.9),
             "This Week.",
             font_size=34, bold=True, color=TEXT)
    add_text(s, Inches(0.8), Inches(1.95), Inches(12.0), Inches(0.5),
             "What to talk to your book about this week — generated from real activity, not guesswork.",
             font_size=14, color=TEXT_SECONDARY, italic=True)

    # The This Week feed card
    card_left = Inches(2.0)
    card_top = Inches(2.8)
    card_w = Inches(9.333)
    card_h = Inches(3.7)

    add_rect(s, card_left, card_top, card_w, card_h,
             fill=WHITE, line_color=BORDER, line_width=1)
    # Header band
    add_rect(s, card_left, card_top, card_w, Inches(0.7),
             fill=SURFACE, line_color=BORDER, line_width=0.5)
    add_text(s, card_left + Inches(0.3), card_top + Inches(0.13),
             Inches(4.0), Inches(0.3),
             "THIS WEEK", font_size=10, bold=True, color=TEXT_MUTED)
    add_text(s, card_left + Inches(0.3), card_top + Inches(0.36),
             Inches(5.0), Inches(0.3),
             "4 clients to reach out to", font_size=15, bold=True, color=TEXT)
    add_text(s, card_left + Inches(5.0), card_top + Inches(0.25),
             card_w - Inches(5.3), Inches(0.3),
             "Mon, May 11", font_size=10, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)

    # Rows
    rows = [
        (ACCENT_RED, "⏰", "Sarah Westlake", "Auto renewal · Allstate · 21 days out · premium up 12%"),
        (ACCENT_AMBER, "📄", "Robert Thompson", "Uploaded new umbrella declarations — needs review"),
        (ACCENT_CYAN, "👁", "Elena Rodriguez", "Viewed her shared renewal review · ready to discuss"),
        (ACCENT_SLATE, "💬", "Marcus Williams", "No interaction in 127 days · still active book"),
    ]
    row_y = card_top + Inches(0.75)
    row_h = Inches(0.72)
    for i, (sev_color, icon, name, detail) in enumerate(rows):
        y = row_y + Inches(i * 0.72)
        # bottom border
        if i < 3:
            add_rect(s, card_left, y + row_h - Emu(1), card_w, Emu(1500),
                     fill=BORDER)
        # Severity dot + icon
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  card_left + Inches(0.3),
                                  y + Inches(0.2),
                                  Inches(0.32), Inches(0.32))
        dot.fill.solid()
        dot.fill.fore_color.rgb = sev_color
        dot.line.fill.background()
        add_text(s, card_left + Inches(0.3), y + Inches(0.2),
                 Inches(0.32), Inches(0.32),
                 icon, font_size=12, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, bold=True)
        # Name + detail
        add_text(s, card_left + Inches(0.85), y + Inches(0.12),
                 card_w - Inches(2.2), Inches(0.3),
                 name, font_size=12, bold=True, color=TEXT)
        add_text(s, card_left + Inches(0.85), y + Inches(0.4),
                 card_w - Inches(2.2), Inches(0.3),
                 detail, font_size=10, color=TEXT_MUTED)
        # Open link
        add_text(s, card_left + card_w - Inches(0.9), y + Inches(0.25),
                 Inches(0.7), Inches(0.3),
                 "Open →", font_size=10, bold=True, color=PRIMARY,
                 align=PP_ALIGN.RIGHT)

    add_text(s, Inches(0.8), Inches(6.65), Inches(12.0), Inches(0.3),
             "Every row is an observation from real activity in the book — no scores, no predictions, no \"AI thinks you should call X.\"",
             font_size=10, color=TEXT_MUTED, italic=True, align=PP_ALIGN.CENTER)

    page_footer(s, 5)


# ── Slide 6: Where AI does the work ─────────────────────────────────────────

def slide_6(prs: Presentation):
    s = add_blank_slide(prs)
    set_slide_bg(s, WHITE)

    add_text(s, Inches(0.8), Inches(0.6), Inches(12.0), Inches(0.4),
             "HOW IT WORKS", font_size=10, bold=True, color=PRIMARY)
    add_text(s, Inches(0.8), Inches(1.05), Inches(12.0), Inches(1.0),
             "AI that helps. Never AI that decides.",
             font_size=28, bold=True, color=TEXT)
    add_text(s, Inches(0.8), Inches(2.0), Inches(12.0), Inches(0.5),
             "Covrabl uses AI to simplify, surface, and speed up the work — so your time goes to the conversations only you can have.",
             font_size=14, color=TEXT_SECONDARY)

    # 5 capability tiles in a horizontal row
    capabilities = [
        ("📖", "Reads policies",
         "Extracts coverage data from any PDF, scan, or photo."),
        ("🗣", "Translates jargon",
         "Plain-English summaries clients actually understand."),
        ("🔁", "Compares year-over-year",
         "Surfaces what changed without re-reading 90 pages."),
        ("✓", "Verifies compliance",
         "Matches actual coverage against lease, loan, or vendor requirements."),
        ("👀", "Watches the book",
         "Surfaces who needs attention this week from real activity."),
    ]
    tile_top = Inches(3.7)
    tile_h = Inches(2.4)
    tile_w = Inches(2.35)
    gap = Inches(0.1)
    left0 = Inches(0.6)

    for i, (icon, title, body) in enumerate(capabilities):
        left = left0 + i * (tile_w + gap)
        add_rect(s, left, tile_top, tile_w, tile_h, fill=SURFACE,
                 line_color=BORDER, line_width=1)
        add_text(s, left + Inches(0.2), tile_top + Inches(0.25),
                 tile_w - Inches(0.4), Inches(0.6),
                 icon, font_size=26, color=PRIMARY, align=PP_ALIGN.CENTER)
        add_text(s, left + Inches(0.2), tile_top + Inches(0.95),
                 tile_w - Inches(0.4), Inches(0.5),
                 title, font_size=12, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
        add_text(s, left + Inches(0.2), tile_top + Inches(1.45),
                 tile_w - Inches(0.4), Inches(0.85),
                 body, font_size=10, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)

    # Footer guardrail line
    add_text(s, Inches(0.8), Inches(6.5), Inches(12.0), Inches(0.4),
             ("Every AI output is observable and overridable. "
              "Covrabl shows its work; the agent decides what matters."),
             font_size=11, italic=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    page_footer(s, 6)


# ── Slide 7: How We Fit ─────────────────────────────────────────────────────

def slide_7(prs: Presentation):
    s = add_blank_slide(prs)
    set_slide_bg(s, WHITE)

    add_text(s, Inches(0.8), Inches(0.6), Inches(12.0), Inches(0.4),
             "WHERE COVRABL SITS", font_size=10, bold=True, color=PRIMARY)
    add_text(s, Inches(0.8), Inches(1.05), Inches(12.0), Inches(1.6),
             "We don't replace your AMS. We layer above it as the client-facing engagement surface.",
             font_size=26, bold=True, color=TEXT)

    # Three boxes left-to-right: AMS → Covrabl → Client
    box_top = Inches(3.4)
    box_h = Inches(1.7)
    box_w = Inches(3.5)
    gap = Inches(0.4)
    left1 = Inches(0.8)
    left2 = left1 + box_w + gap
    left3 = left2 + box_w + gap

    boxes = [
        (left1, "AMS / CRM",
         "Stores policies. Runs servicing. Handles commissions.", SURFACE, TEXT, BORDER),
        (left2, "COVRABL",
         "Reads your book. Surfaces engagement signals. Presents to clients.",
         PRIMARY, WHITE, PRIMARY),
        (left3, "CLIENT",
         "Sees their coverage. Reads your reviews. Engages between renewals.",
         SURFACE, TEXT, BORDER),
    ]
    for left, title, body, fill, fg, border in boxes:
        add_rect(s, left, box_top, box_w, box_h, fill=fill, line_color=border, line_width=1)
        add_text(s, left + Inches(0.3), box_top + Inches(0.3),
                 box_w - Inches(0.6), Inches(0.4),
                 title, font_size=14, bold=True, color=fg)
        add_text(s, left + Inches(0.3), box_top + Inches(0.85),
                 box_w - Inches(0.6), Inches(0.8),
                 body, font_size=11, color=fg if fill != PRIMARY else WHITE)

    # Arrows between boxes
    for x in [left2 - Inches(0.3), left3 - Inches(0.3)]:
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                  x, box_top + Inches(0.7), Inches(0.3), Inches(0.3))
        arr.fill.solid()
        arr.fill.fore_color.rgb = TEXT_MUTED
        arr.line.fill.background()

    # Onboarding line
    add_text(s, Inches(0.8), Inches(5.7), Inches(12.0), Inches(0.4),
             "Onboarding: 15 minutes to invite your first client. Bulk-add via CSV for larger books.",
             font_size=14, color=TEXT_SECONDARY, italic=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.8), Inches(6.15), Inches(12.0), Inches(0.4),
             "Your producers. Your branding. Your renewal calendar — already familiar.",
             font_size=12, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    page_footer(s, 7)


# ── Slide 8: Why This Drives Retention ──────────────────────────────────────

def slide_8(prs: Presentation):
    s = add_blank_slide(prs)
    set_slide_bg(s, WHITE)

    add_text(s, Inches(0.8), Inches(0.6), Inches(12.0), Inches(0.4),
             "THE DEFENSIVE WIN", font_size=10, bold=True, color=PRIMARY)
    add_text(s, Inches(0.8), Inches(1.05), Inches(12.0), Inches(1.6),
             "Agencies that talk to clients more than once a year don't get shopped at renewal.",
             font_size=26, bold=True, color=TEXT)

    # Two columns
    col_top = Inches(3.4)
    col_h = Inches(3.0)
    col_w = Inches(5.9)
    gap = Inches(0.5)
    left1 = Inches(0.8)
    left2 = left1 + col_w + gap

    # Left — retention benchmark
    add_rect(s, left1, col_top, col_w, col_h, fill=SURFACE,
             line_color=BORDER, line_width=1)
    add_text(s, left1 + Inches(0.4), col_top + Inches(0.3),
             col_w - Inches(0.8), Inches(0.4),
             "INDUSTRY RETENTION BENCHMARK", font_size=10, bold=True, color=TEXT_MUTED)
    add_text(s, left1 + Inches(0.4), col_top + Inches(0.8),
             col_w - Inches(0.8), Inches(0.7),
             "Top-quartile agencies: 93%+",
             font_size=20, bold=True, color=TEXT)
    add_text(s, left1 + Inches(0.4), col_top + Inches(1.5),
             col_w - Inches(0.8), Inches(0.6),
             "Median agencies: ~88%",
             font_size=20, bold=True, color=TEXT_SECONDARY)
    add_text(s, left1 + Inches(0.4), col_top + Inches(2.2),
             col_w - Inches(0.8), Inches(0.4),
             "The five-point gap is where retention is won.",
             font_size=12, italic=True, color=TEXT_SECONDARY)
    add_text(s, left1 + Inches(0.4), col_top + Inches(2.55),
             col_w - Inches(0.8), Inches(0.3),
             "Source: IIABA Best Practices study.",
             font_size=9, color=TEXT_MUTED, italic=True)

    # Right — what Covrabl unlocks
    add_rect(s, left2, col_top, col_w, col_h, fill=WHITE,
             line_color=PRIMARY, line_width=1.5)
    add_text(s, left2 + Inches(0.4), col_top + Inches(0.3),
             col_w - Inches(0.8), Inches(0.4),
             "WHAT COVRABL UNLOCKS", font_size=10, bold=True, color=PRIMARY)
    bullets = [
        "Structured renewal review every client gets",
        "Between-renewal outreach feed (This Week)",
        "Branded client portal as the always-on touchpoint",
    ]
    by = col_top + Inches(0.9)
    for i, b in enumerate(bullets):
        add_text(s, left2 + Inches(0.4), by + Inches(i * 0.65),
                 col_w - Inches(0.8), Inches(0.6),
                 f"—  {b}", font_size=14, color=TEXT)

    add_text(s, left2 + Inches(0.4), col_top + col_h - Inches(0.5),
             col_w - Inches(0.8), Inches(0.4),
             "Less work. More contact. Higher retention.",
             font_size=12, italic=True, color=PRIMARY, bold=True)

    page_footer(s, 8)


# ── Slide 9: Trust & Security ───────────────────────────────────────────────

def slide_9(prs: Presentation):
    s = add_blank_slide(prs)
    set_slide_bg(s, WHITE)

    add_text(s, Inches(0.8), Inches(0.6), Inches(12.0), Inches(0.4),
             "AGENCY TRUST INFRASTRUCTURE", font_size=10, bold=True, color=PRIMARY)
    add_text(s, Inches(0.8), Inches(1.05), Inches(12.0), Inches(1.1),
             "Built for the trust your book has in you.",
             font_size=28, bold=True, color=TEXT)
    add_text(s, Inches(0.8), Inches(2.15), Inches(12.0), Inches(0.6),
             ("When you invite your clients into Covrabl, you're extending your agency's "
              "reputation. We treat that as the responsibility it is."),
             font_size=14, color=TEXT_SECONDARY)

    # Four trust pillars
    pillars = [
        ("No data sold. Ever.",
         "We make money from agencies, not advertisers or lead aggregators."),
        ("No carrier marketplace.",
         "Covrabl will never quote against your business or resell client data."),
        ("Assistive AI, never authoritative.",
         "Covrabl surfaces observations; the agent stays the authority."),
        ("Encrypted at rest. MFA. Audit log.",
         "AES-256. Two-factor auth shipped. Audit log on every account."),
    ]
    p_top = Inches(3.3)
    p_h = Inches(1.5)
    p_w = Inches(2.95)
    gap = Inches(0.15)
    left0 = Inches(0.8)
    for i, (head, body) in enumerate(pillars):
        left = left0 + i * (p_w + gap)
        add_rect(s, left, p_top, p_w, p_h, fill=SURFACE,
                 line_color=BORDER, line_width=1)
        add_text(s, left + Inches(0.3), p_top + Inches(0.3),
                 p_w - Inches(0.6), Inches(0.6),
                 head, font_size=13, bold=True, color=TEXT)
        add_text(s, left + Inches(0.3), p_top + Inches(0.85),
                 p_w - Inches(0.6), Inches(0.6),
                 body, font_size=10, color=TEXT_SECONDARY)

    # Footer
    add_text(s, Inches(0.8), Inches(5.5), Inches(12.0), Inches(0.4),
             "SOC 2 readiness in progress. Subprocessor list and privacy policy public at covrabl.com.",
             font_size=11, italic=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    page_footer(s, 9)


# ── Slide 10: Founding Partner ──────────────────────────────────────────────

def slide_10(prs: Presentation):
    s = add_blank_slide(prs)
    set_slide_bg(s, WHITE)

    add_text(s, Inches(0.8), Inches(0.6), Inches(12.0), Inches(0.4),
             "FOUNDING PARTNER PRICING", font_size=10, bold=True, color=PRIMARY)
    add_text(s, Inches(0.8), Inches(1.05), Inches(12.0), Inches(1.0),
             "Locked for the first 12 months.",
             font_size=28, bold=True, color=TEXT)

    # Vision footnote — small italic, above the pricing block
    add_text(s, Inches(0.8), Inches(2.15), Inches(12.0), Inches(0.5),
             ("Where this goes: over time, Covrabl becomes the visibility and "
              "verification layer across the client relationship. The wedge is "
              "between-renewal engagement. The longer game is broader."),
             font_size=11, italic=True, color=TEXT_MUTED)

    # Divider rule
    add_rect(s, Inches(0.8), Inches(2.85), Inches(11.733), Emu(1500), fill=BORDER)

    # Two pricing cards
    card_top = Inches(3.1)
    card_h = Inches(3.0)
    card_w = Inches(5.7)
    gap = Inches(0.3)
    left1 = Inches(0.8)
    left2 = left1 + card_w + gap

    # Founding partner card (highlighted)
    add_rect(s, left1, card_top, card_w, card_h, fill=WHITE,
             line_color=PRIMARY, line_width=2.5)
    add_text(s, left1 + Inches(0.4), card_top + Inches(0.3),
             card_w - Inches(0.8), Inches(0.3),
             "FOUNDING PARTNER", font_size=10, bold=True, color=PRIMARY)
    add_text(s, left1 + Inches(0.4), card_top + Inches(0.7),
             card_w - Inches(0.8), Inches(0.9),
             "$59 / month per agent",
             font_size=28, bold=True, color=TEXT)
    add_text(s, left1 + Inches(0.4), card_top + Inches(1.55),
             card_w - Inches(0.8), Inches(0.4),
             "Locks for 12 months · Every feature · Unlimited clients",
             font_size=11, color=TEXT_SECONDARY)
    feats = [
        "Coverage Reviews (renewal + quote)",
        "This Week outreach feed",
        "Branded client portal",
        "Audit log + data export",
    ]
    fy = card_top + Inches(2.0)
    for i, f in enumerate(feats):
        add_text(s, left1 + Inches(0.4), fy + Inches(i * 0.22),
                 card_w - Inches(0.8), Inches(0.25),
                 f"✓  {f}", font_size=10, color=TEXT)

    # Add-on card
    add_rect(s, left2, card_top, card_w, card_h, fill=SURFACE,
             line_color=BORDER, line_width=1)
    add_text(s, left2 + Inches(0.4), card_top + Inches(0.3),
             card_w - Inches(0.8), Inches(0.3),
             "ADD-ON", font_size=10, bold=True, color=TEXT_MUTED)
    add_text(s, left2 + Inches(0.4), card_top + Inches(0.7),
             card_w - Inches(0.8), Inches(0.6),
             "White-label",
             font_size=22, bold=True, color=TEXT)
    add_text(s, left2 + Inches(0.4), card_top + Inches(1.3),
             card_w - Inches(0.8), Inches(0.4),
             "$500 one-time setup",
             font_size=14, color=TEXT_SECONDARY)
    add_text(s, left2 + Inches(0.4), card_top + Inches(1.85),
             card_w - Inches(0.8), Inches(1.0),
             ("Your logo, your colors, your domain on the client-facing portal. "
              "Your clients never see a Covrabl brand."),
             font_size=11, color=TEXT_SECONDARY)

    # Below cards
    add_text(s, Inches(0.8), Inches(6.3), Inches(12.0), Inches(0.3),
             "After founding-partner closes: $99/mo list.",
             font_size=11, italic=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.8), Inches(6.65), Inches(12.0), Inches(0.4),
             f"Book a 15-min demo with the founder  ·  {CONTACT_EMAIL}  ·  {WEBSITE}",
             font_size=13, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

    page_footer(s, 10)


# ── Slide 11: Closing ───────────────────────────────────────────────────────

def slide_11(prs: Presentation):
    s = add_blank_slide(prs)
    set_slide_bg(s, NAVY_DARK)

    # Big block of text, centered
    paragraphs = [
        ("Insurance relationships shouldn't disappear between renewals.", 26, False, RGBColor(0xCB, 0xD5, 0xE1)),
        ("", 8, False, WHITE),
        ("Renewal happens. The policy lands. The conversation ends.", 18, False, RGBColor(0x94, 0xA3, 0xB8)),
        ("", 8, False, WHITE),
        ("Until something breaks.", 18, True, WHITE),
        ("", 16, False, WHITE),
        ("Covrabl is what lives in the silence.", 30, True, SECONDARY),
    ]
    tb = s.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(11.333), Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, (text, size, bold, color) in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

    # Bottom-right CTA chip
    add_text(s, Inches(7.0), Inches(6.7), Inches(5.5), Inches(0.4),
             f"{CONTACT_EMAIL}  ·  {WEBSITE}",
             font_size=11, color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.RIGHT)


# ── Build ───────────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_1(prs)
    slide_2(prs)
    slide_3(prs)
    slide_4(prs)
    slide_5(prs)
    slide_6(prs)
    slide_7(prs)
    slide_8(prs)
    slide_9(prs)
    slide_10(prs)
    slide_11(prs)

    # If the primary file is locked (e.g. you have PowerPoint open),
    # fall back to a versioned filename so the regen never fails silently.
    saved_any = False
    for out in OUTPUT_PATHS or [OUTPUT_PATH]:
        os.makedirs(out.parent, exist_ok=True)
        try:
            prs.save(out)
            print(f"Saved {out}")
            saved_any = True
        except PermissionError:
            alt = out.with_name(out.stem + "-NEW" + out.suffix)
            prs.save(alt)
            print(f"Primary locked; saved {alt}  (close PowerPoint to refresh the original)")
            saved_any = True
    if not saved_any:
        raise RuntimeError("Could not save deck to any desktop path")


if __name__ == "__main__":
    build()
